"""
Material parsing API view.

Endpoint:
    POST /api/v1/courses/ai/parse-material/  -- Extract text from uploaded files
"""

import logging
import os
from typing import Any

from django.core.cache import cache
from rest_framework import status
from rest_framework.decorators import api_view, parser_classes, permission_classes
from rest_framework.parsers import MultiPartParser
from rest_framework.permissions import IsAuthenticated
from rest_framework.response import Response

from utils.decorators import admin_only, tenant_required

logger = logging.getLogger(__name__)

# Rate limit: 10 parse requests per hour per tenant
PARSE_RATE_LIMIT = 10
PARSE_RATE_WINDOW = 3600  # seconds

# Max file size: 50 MB
MAX_FILE_SIZE = 50 * 1024 * 1024

# Max extracted text length
MAX_TEXT_LENGTH = 100_000

# Allowed MIME types grouped by category
ALLOWED_MIME_TYPES = {
    "pdf": {"application/pdf"},
    "docx": {
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    },
    "pptx": {
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    },
    "video": {
        "video/mp4",
        "video/mpeg",
        "video/quicktime",
        "video/x-msvideo",
        "video/webm",
        "video/x-matroska",
        "video/ogg",
    },
}

# Flat set of all allowed MIME types
ALL_ALLOWED_MIMES = set()
for _mimes in ALLOWED_MIME_TYPES.values():
    ALL_ALLOWED_MIMES.update(_mimes)


def _check_parse_rate_limit(tenant_id: str) -> bool:
    """
    Check if the tenant has exceeded the material parsing rate limit.
    Returns True if the request is allowed, False if rate-limited.
    """
    cache_key = f"ai_parse_material_rate:{tenant_id}"
    cache.add(cache_key, 0, timeout=PARSE_RATE_WINDOW)
    current_count = cache.incr(cache_key)
    return current_count <= PARSE_RATE_LIMIT


def _detect_file_category(content_type: str, file_name: str) -> str | None:
    """
    Detect the file category from MIME type, with extension-based fallback.
    Returns one of: 'pdf', 'docx', 'pptx', 'video', or None if unsupported.
    """
    # Check MIME type first
    for category, mimes in ALLOWED_MIME_TYPES.items():
        if content_type in mimes:
            return category

    # Also accept any video/* MIME
    if content_type.startswith("video/"):
        return "video"

    # Extension-based fallback for common cases where MIME detection fails
    ext = os.path.splitext(file_name)[1].lower()
    ext_map = {
        ".pdf": "pdf",
        ".docx": "docx",
        ".pptx": "pptx",
        ".mp4": "video",
        ".mov": "video",
        ".avi": "video",
        ".webm": "video",
        ".mkv": "video",
    }
    return ext_map.get(ext)


def _parse_pdf(file_obj) -> dict[str, Any]:
    """
    Extract text from a PDF file.
    Returns: { text, metadata: { pages, ... } }
    """
    try:
        import PyPDF2
    except ImportError:
        logger.error("PyPDF2 is not installed. Install with: pip install PyPDF2")
        return {
            "text": "",
            "metadata": {"pages": 0},
            "error": "PDF parsing library not available. Please install PyPDF2.",
        }

    try:
        reader = PyPDF2.PdfReader(file_obj)
        num_pages = len(reader.pages)
        text_parts: list[str] = []
        total_len = 0

        for page_num, page in enumerate(reader.pages):
            try:
                page_text = page.extract_text() or ""
                if page_text:
                    text_parts.append(page_text)
                    total_len += len(page_text)
                    if total_len >= MAX_TEXT_LENGTH:
                        break
            except Exception as e:
                logger.warning("Failed to extract text from PDF page %d: %s", page_num + 1, e)
                continue

        full_text = "\n\n".join(text_parts)[:MAX_TEXT_LENGTH]

        return {
            "text": full_text,
            "metadata": {"pages": num_pages},
        }
    except Exception as e:
        logger.warning("PDF parsing failed: %s", e)
        return {
            "text": "",
            "metadata": {"pages": 0},
            "error": f"Failed to parse PDF: {str(e)}",
        }


def _parse_docx(file_obj) -> dict[str, Any]:
    """
    Extract text from a DOCX file.
    Returns: { text, metadata: { ... } }
    """
    try:
        import docx
    except ImportError:
        logger.error("python-docx is not installed. Install with: pip install python-docx")
        return {
            "text": "",
            "metadata": {},
            "error": "DOCX parsing library not available. Please install python-docx.",
        }

    try:
        document = docx.Document(file_obj)
        text_parts: list[str] = []

        # Extract paragraph text
        for para in document.paragraphs:
            text = para.text.strip()
            if text:
                text_parts.append(text)

        # Extract table text
        for table in document.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text_parts.append(row_text)

        full_text = "\n\n".join(text_parts)[:MAX_TEXT_LENGTH]

        return {
            "text": full_text,
            "metadata": {
                "paragraphs": len(document.paragraphs),
                "tables": len(document.tables),
            },
        }
    except Exception as e:
        logger.warning("DOCX parsing failed: %s", e)
        return {
            "text": "",
            "metadata": {},
            "error": f"Failed to parse DOCX: {str(e)}",
        }


def _parse_pptx(file_obj) -> dict[str, Any]:
    """
    Extract text from a PPTX file.
    Returns: { text, metadata: { slides, ... } }
    """
    try:
        from pptx import Presentation
    except ImportError:
        logger.error("python-pptx is not installed. Install with: pip install python-pptx")
        return {
            "text": "",
            "metadata": {"slides": 0},
            "error": "PPTX parsing library not available. Please install python-pptx.",
        }

    try:
        prs = Presentation(file_obj)
        text_parts: list[str] = []
        num_slides = len(prs.slides)

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_texts: list[str] = []

            # Extract text from shapes
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)

                # Extract text from tables in shapes
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            slide_texts.append(row_text)

            # Extract notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_texts.append(f"[Speaker Notes: {notes_text}]")

            if slide_texts:
                text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_texts))

        full_text = "\n\n".join(text_parts)[:MAX_TEXT_LENGTH]

        return {
            "text": full_text,
            "metadata": {"slides": num_slides},
        }
    except Exception as e:
        logger.warning("PPTX parsing failed: %s", e)
        return {
            "text": "",
            "metadata": {"slides": 0},
            "error": f"Failed to parse PPTX: {str(e)}",
        }


def _handle_video(file_obj, file_name: str, file_size: int) -> dict[str, Any]:
    """
    Handle video file upload. Returns a stub response since video
    transcription is handled asynchronously by the existing Whisper pipeline.
    """
    return {
        "text": "",
        "metadata": {},
        "processing": "async",
        "message": "Video transcription will be processed asynchronously",
    }


@api_view(["POST"])
@parser_classes([MultiPartParser])
@permission_classes([IsAuthenticated])
@admin_only
@tenant_required
def parse_material(request):
    """
    Parse an uploaded file and extract text content.

    Accepts: multipart/form-data with a 'file' field
    Supported formats: PDF, DOCX, PPTX, video/*

    Returns:
        200: { text, metadata: { pages?, file_type, file_size, file_name } }
        400: Validation error
        429: Rate limit exceeded
    """
    tenant_id = str(request.tenant.id)

    if not _check_parse_rate_limit(tenant_id):
        return Response(
            {
                "error": "Rate limit exceeded. You can parse up to 10 files per hour.",
                "retry_after_seconds": PARSE_RATE_WINDOW,
            },
            status=status.HTTP_429_TOO_MANY_REQUESTS,
        )

    uploaded_file = request.FILES.get("file")
    if not uploaded_file:
        return Response(
            {"error": "No file provided. Include a 'file' field in the request."},
            status=status.HTTP_400_BAD_REQUEST,
        )

    file_name = uploaded_file.name or "unknown"
    file_size = uploaded_file.size
    content_type = uploaded_file.content_type or ""

    # Validate file size
    if file_size > MAX_FILE_SIZE:
        max_mb = MAX_FILE_SIZE // (1024 * 1024)
        return Response(
            {"error": f"File too large. Maximum allowed size is {max_mb}MB."},
            status=status.HTTP_400_BAD_REQUEST,
        )

    # Detect file category
    file_category = _detect_file_category(content_type, file_name)
    if file_category is None:
        return Response(
            {
                "error": (
                    f"Unsupported file type: {content_type or 'unknown'}. "
                    "Supported formats: PDF, DOCX, PPTX, and video files."
                )
            },
            status=status.HTTP_400_BAD_REQUEST,
        )

    logger.info(
        "parse_material: processing %s (%s, %d bytes) for tenant=%s",
        file_name,
        file_category,
        file_size,
        tenant_id,
    )

    # Parse based on file type
    if file_category == "pdf":
        result = _parse_pdf(uploaded_file)
    elif file_category == "docx":
        result = _parse_docx(uploaded_file)
    elif file_category == "pptx":
        result = _parse_pptx(uploaded_file)
    elif file_category == "video":
        result = _handle_video(uploaded_file, file_name, file_size)
    else:
        return Response(
            {"error": f"Unsupported file category: {file_category}"},
            status=status.HTTP_400_BAD_REQUEST,
        )

    # Build response metadata
    base_metadata = {
        "file_type": file_category,
        "file_size": file_size,
        "file_name": file_name,
    }
    result_metadata = result.get("metadata", {})
    base_metadata.update(result_metadata)

    response_data: dict[str, Any] = {
        "text": result.get("text", ""),
        "metadata": base_metadata,
    }

    # Include processing info for video
    if "processing" in result:
        response_data["processing"] = result["processing"]
        response_data["message"] = result.get("message", "")

    # Include parsing error if partial failure
    if "error" in result:
        response_data["warning"] = result["error"]

    logger.info(
        "parse_material: extracted %d chars from %s (%s)",
        len(response_data["text"]),
        file_name,
        file_category,
    )

    return Response(response_data, status=status.HTTP_200_OK)
